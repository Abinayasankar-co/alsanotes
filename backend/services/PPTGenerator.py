import json5
import ast
from io import BytesIO
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from services.ImageProcessing import StabilityAIClient
from pptx.dml.color import RGBColor
from fastapi.exceptions import HTTPException

class PowerPointGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.image = StabilityAIClient()

    def create_presentation(self,data):
        # Create title slide
        try:
         #self.datas = json.loads(data)
         #self.datas = ast.literal_eval(data)
         self.datas = json5.loads(data)
         print(type(self.datas))
         #print(self.datas)
         self._add_title_slide(self.datas["title"])

         #adding Background 
         self.set_background()
        
         # Create content slides
         for slide_data in self.datas["slides"]:
             self._add_content_slide(slide_data)

         # Save the presentation 
         self.prs.save('./business_performance_presentation.pptx')
         print("PowerPoint with table, bar graph, and pie chart generated successfully!")
        except Exception as e:
            raise HTTPException(status_code=400, detail="Error generating PowerPoint: " + str(e))

            
    def set_background(self):
        """Set a custom background color for a slide."""
        r, g, b = self.datas["colors"]["background"].values()
        background = self.prs.slides.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
    
    def add_image(self,slidedata):
        """
        Add an image to the slide with specified position and size.
        
        Parameters:
        slide: The slide object to which the image will be added.
        image_path: Path of the image file.
        left: Distance from the left side of the slide (Inches).
        top: Distance from the top side of the slide (Inches).
        width: Width of the image (Inches).
        height: Height of the image (Inches). If None, it maintains the aspect ratio.

        """
        try:  
            left , top , width , height = slidedata["imgdata"]["position"].values()
            img_stream = self.image.generate_image(slidedata["imgdata"]["imgprompt"])
            self.prs.shapes.add_picture(BytesIO(img_stream.read()), left, top, width=width, height=height)
        except Exception as e:
            raise HTTPException(status_code=404,detail={f"Error While Generating images ans slides : {e}"})

    def _add_title_slide(self, title_text):
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = title_text

    def _add_content_slide(self, slide_data):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_data["title"]

        # Add content or table
        if "content" in slide_data:
            for content_item in slide_data["content"]:
                content_box = slide.shapes.placeholders[1].text_frame
                content_box.text = content_item["text"]

        if "table_data" in slide_data:
            self._add_table(slide, slide_data["table_data"]["headers"], slide_data["table_data"]["rows"])

        if "chart" in slide_data:
            chart_type = slide_data["chart_type"]
            if chart_type == "bar":
                self._add_bar_chart(slide, slide_data)
            elif chart_type == "pie":
                self._add_pie_chart(slide, slide_data)

        if "imgdata" in slide_data:
            self.add_image(slide_data)

    def _add_table(self, slide, headers, rows):
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(2)
        table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, cx, cy).table

        # Add headers
        for i, header in enumerate(headers):
            table.cell(0, i).text = header

        # Add rows
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, cell_value in enumerate(row):
                table.cell(row_idx, col_idx).text = str(cell_value)

    def _add_bar_chart(self, slide, slide_data):
        image_path = "bar_chart.png"
        categories = list(slide_data["data"].keys())
        values = list(slide_data["data"].values())
        self._create_bar_chart(categories, values, slide_data["title"], image_path)
        slide.shapes.add_picture(image_path, Inches(2), Inches(3), Inches(6), Inches(3))

    def _add_pie_chart(self, slide, slide_data):
        image_path = "pie_chart.png"
        categories = list(slide_data["data"].keys())
        values = list(slide_data["data"].values())
        print(values)
        self._create_pie_chart(categories, values, slide_data["title"], image_path)
        slide.shapes.add_picture(image_path, Inches(2), Inches(3), Inches(6), Inches(3))

    def _create_bar_chart(self, categories, values, title, image_path):
        plt.figure(figsize=(5, 4))
        plt.bar(categories, values, color='blue')
        plt.xlabel('Category')
        plt.ylabel('Values')
        plt.title(title)
        plt.savefig(image_path)
        plt.close()

    def _create_pie_chart(self, categories, values, title, image_path):
        plt.figure(figsize=(5, 4))
        plt.pie(values, labels=categories, autopct='%1.1f%%', startangle=140)
        plt.axis('equal')
        plt.title(title)
        plt.savefig(image_path)
        plt.close()




#Example Usage
"""
json_data = """

{
    "title": "The Himalayas: A Natural and Cultural Wonder",
    "colors": {
       "background": {
           "r": 192,
           "g": 221,
           "b": 255
       }
    },
    "slides": [
       {
           "title": "Introduction to the Himalayas",
           "content": [
               {
                   "text": "The Himalayas, spanning approximately 2,400 kilometers across five countries — India, Nepal, Bhutan, China, and Pakistan — are home to some of the world's highest peaks, including Mount Everest, which stands at 8,848 meters. This mountain range contains over 100 peaks rising above 7,200 meters, making it the tallest and one of the youngest mountain ranges on Earth, formed about 50 million years ago from the collision of the Indian and Eurasian tectonic plates."
               }
           ],
           "imgdata": {
               "imgprompt": "Aerial view of the Himalayas with snow-capped peaks and valleys",
               "position": {
                   "left": 10,
                   "top": 10,
                   "width": 800,
                   "height": 600
               }
       }
       },
       {
           "title": "Biodiversity of the Himalayas",
           "content": [
               {
                   "text": "The Himalayas host some of the rarest and most endangered animals, such as the snow leopard, red panda, and Himalayan musk deer. Forests in the lower elevations consist of subtropical pine and temperate broadleaf forests, while the higher altitudes are covered by alpine shrubs and meadows."
               }
           ],
           "imgdata": {
               "imgprompt": "Image of a snow leopard in the Himalayas",
               "position": {
                   "left": 10,
                   "top": 10,
                   "width": 800,
                   "height": 600
               }
       }
       },
       {
           "title": "Freshwater Resources of the Himalayas",
           "content": [
               {
                   "text": "The Himalayas are a critical source of freshwater, with glaciers and snowfields supplying about 8.6 million cubic meters of water annually to major river systems like the Ganges, Indus, and Brahmaputra. These rivers support nearly 1.5 billion people across the South Asian subcontinent, providing water for agriculture, drinking, and hydroelectric power generation."
               }
           ],
           "imgdata": {
               "imgprompt": "Image of a glacier in the Himalayas",
               "position": {
                   "left": 10,
                   "top": 10,
                   "width": 800,
                   "height": 600
               }
       }
       },
       {
           "title": "Cultural Significance of the Himalayas",
           "content": [
               {
                   "text": "The Himalayas have been a spiritual and cultural center for centuries, with numerous temples, monasteries, and pilgrimage sites scattered throughout the region. The mountain range has also inspired countless works of art, literature, and music, reflecting its profound impact on the human imagination."
               }
           ],
           "imgdata": {
               "imgprompt": "Image of a Buddhist monastery in the Himalayas",
               "position": {
                   "left": 10,
                   "top": 10,
                   "width": 800,
                   "height": 600
               }
       }
       },
       {
           "title": "Challenges Facing the Himalayas",
           "content": [
               {
                   "text": "The Himalayas face numerous challenges, including climate change, deforestation, and overtourism. These threats not only endanger the region's fragile ecosystems but also threaten the livelihoods of the millions of people who depend on the Himalayas for their survival."
               }
           ],
           "imgdata": {
               "imgprompt": "Image of deforestation in the Himalayas",
               "position": {
                   "left": 10,
                   "top": 10,
                   "width": 800,
                   "height": 600
               }
       }
       },
       {
           "title": "Conservation Efforts in the Himalayas",
           "content": [
               {
                   "text": "Numerous organizations and governments are working to protect the Himalayas and its unique biodiversity. These efforts include habitat restoration, community-based conservation, and sustainable tourism initiatives that aim to balance the needs of people and the environment.",
               }
           ],
           "imgdata": {
               "imgprompt": "Image of a conservation project in the Himalayas",
               "position": {
                   "left": 10,
                   "top": 10,
                   "width": 800,
                   "height": 600
               }
       }
       },
       {
           "title": "The Future of the Himalayas",
           "content": [
               {
                   "text": "The Himalayas are at a critical juncture, with the potential for both destruction and renewal. By working together, we can ensure that this natural and cultural wonder remains a source of inspiration and sustenance for generations to come."
               }
           ],
           "imgdata": {
               "imgprompt": "Image of a sunrise over the Himalayas",
               "position": {
                   "left": 10,
                   "top": 10,
                   "width": 800,
                   "height": 600
               }
       }
    }
    ]
}

"""
ppt_generator = PowerPointGenerator()
ppt_generator.create_presentation(json_data)


"""

